import os
import datetime
import io
import zipfile
import tkinter as tk
import logging
from tkinter import filedialog, simpledialog, messagebox
import gridfs
import pypdf
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Cargar variables de entorno (centralizado en src.config)
from src.config import DB_NAME, COLS
from src.database import get_database

logger = logging.getLogger(__name__)

# --- 1. CONFIGURACIÓN (ATLAS) ---
COLLECTION_RAW = COLS["RAW"]

def pedir_usuario_gui():
    """Abre un input dialog para pedir el nombre del usuario."""
    try:
        root = tk.Tk()
        root.withdraw() # Ocultar ventana principal
        root.attributes('-topmost', True) # Poner al frente
        
        nombre = simpledialog.askstring("Identificación", "Por favor, ingresa tu nombre de usuario para asociar los archivos:")
        
        root.destroy()
        return nombre
    except Exception as e:
        logger.error(f"⚠️ Error al pedir usuario: {e}")
        # Fallback por consola si falla la GUI
        return input("Ingresa tu nombre de usuario: ")

def seleccionar_archivos_gui(initial_dir=None):
    """Abre una ventana del gestor de archivos para seleccionar archivos."""
    try:
        root = tk.Tk()
        root.withdraw()
        root.attributes('-topmost', True)

        filetypes = [("Documentos", ("*.pdf", "*.docx", "*.pptx")), ("Todos los archivos", "*.*")]
        paths = filedialog.askopenfilenames(
            title='Seleccionar archivo(s) a procesar', 
            initialdir=initial_dir or '.', 
            filetypes=filetypes
        )
        root.destroy()
        return list(paths)
    except Exception as e:
        logger.error(f"⚠️ Error en selector de archivos: {e}")
        return []

def conectar_bd():
    try:
        db = get_database(DB_NAME)
        fs = gridfs.GridFS(db)
        return db[COLLECTION_RAW], fs
    except Exception as e:
        logger.error(f"Error conectando a base de datos: {e}")
        return None, None

# --- 2. FUNCIONES AUXILIARES (GridFS) ---

def guardar_imagen_gridfs(fs, img_bytes, nombre_archivo, pagina_idx, img_idx, ext, usuario):
    """Guarda una imagen en GridFS incluyendo el usuario en los metadatos."""
    try:
        filename = f"{usuario}_{nombre_archivo}_P{pagina_idx}_IMG{img_idx}{ext}"
        file_id = fs.put(
            img_bytes, 
            filename=filename,
            metadata={
                "documento_padre": nombre_archivo,
                "pagina_origen": pagina_idx,
                "usuario_propietario": usuario
            }
        )
        return {
            "gridfs_id": file_id,
            "nombre_archivo": filename,
            "tipo_mime": f"image/{ext.replace('.', '')}"
        }
    except Exception as e:
        logger.error(f"⚠️ Error guardando imagen {filename}: {e}")
        return None

# --- 3. EXTRACTORES POR PÁGINA/DIAPOSITIVA ---

def procesar_pdf(ruta_archivo, fs, usuario):
    nombre_doc = os.path.basename(ruta_archivo)
    paginas_estructuradas = []
    
    try:
        reader = pypdf.PdfReader(ruta_archivo)
        for i, page in enumerate(reader.pages):
            idx_pag = i + 1
            texto_pag = page.extract_text() or ""
            
            # Procesar imágenes
            imagenes_pag = []
            for j, img in enumerate(page.images):
                ext = os.path.splitext(img.name)[1]
                # Pasamos 'usuario' para que las imágenes también tengan esa metadata
                img_info = guardar_imagen_gridfs(fs, img.data, nombre_doc, idx_pag, j+1, ext, usuario)
                if img_info:
                    imagenes_pag.append(img_info)
            
            paginas_estructuradas.append({
                "indice": idx_pag,
                "tipo_unidad": "pagina",
                "contenido_texto": texto_pag,
                "imagenes": imagenes_pag,
                "metadata_bloom": None
            })
                
        return paginas_estructuradas
    except Exception as e:
        logger.error(f"⚠️ Error PDF: {e}")
        return []

def procesar_pptx(ruta_archivo, fs, usuario):
    nombre_doc = os.path.basename(ruta_archivo)
    diapositivas_estructuradas = []
    
    try:
        prs = Presentation(ruta_archivo)
        for i, slide in enumerate(prs.slides):
            idx_slide = i + 1
            texto_slide = ""
            imagenes_slide = []
            
            img_count = 0
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto_slide += shape.text + "\n"
                
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_count += 1
                    image = shape.image
                    ext = f".{image.ext}"
                    img_info = guardar_imagen_gridfs(fs, image.blob, nombre_doc, idx_slide, img_count, ext, usuario)
                    if img_info:
                        imagenes_slide.append(img_info)

            diapositivas_estructuradas.append({
                "indice": idx_slide,
                "tipo_unidad": "diapositiva",
                "contenido_texto": texto_slide,
                "imagenes": imagenes_slide,
                "metadata_bloom": None
            })
            
        return diapositivas_estructuradas
    except Exception as e:
        logger.error(f"⚠️ Error PPTX: {e}")
        return []

def procesar_docx(ruta_archivo, fs, usuario):
    nombre_doc = os.path.basename(ruta_archivo)
    imagenes_globales = []
    
    # Extracción de imágenes via ZIP
    try:
        with zipfile.ZipFile(ruta_archivo) as z:
            img_count = 0
            for file_info in z.infolist():
                if file_info.filename.startswith('word/media/'):
                    img_count += 1
                    img_data = z.read(file_info)
                    ext = os.path.splitext(file_info.filename)[1]
                    img_info = guardar_imagen_gridfs(fs, img_data, nombre_doc, 1, img_count, ext, usuario)
                    if img_info:
                        imagenes_globales.append(img_info)
    except:
        pass

    try:
        doc = Document(ruta_archivo)
        texto_completo = "\n".join([p.text for p in doc.paragraphs])
        
        return [{
            "indice": 1,
            "tipo_unidad": "documento_completo",
            "contenido_texto": texto_completo,
            "imagenes": imagenes_globales,
            "metadata_bloom": None
        }]
    except Exception as e:
        logger.error(f"⚠️ Error DOCX: {e}")
        return []

# --- 4. PROCESO PRINCIPAL DE INGESTA ---

def ingestar_archivo(ruta_archivo, collection, fs, usuario):
    """
    Procesa un archivo y lo guarda en MongoDB asociado al usuario.
    """
    if not os.path.exists(ruta_archivo):
        return

    nombre = os.path.basename(ruta_archivo)
    ext = os.path.splitext(nombre)[1].lower()
    logger.info(f"🔄 Procesando: {nombre} (Usuario: {usuario})...")

    unidades_contenido = []

    # Pasamos 'usuario' a las funciones de procesamiento
    if ext == '.pdf':
        unidades_contenido = procesar_pdf(ruta_archivo, fs, usuario)
    elif ext == '.pptx':
        unidades_contenido = procesar_pptx(ruta_archivo, fs, usuario)
    elif ext == '.docx':
        unidades_contenido = procesar_docx(ruta_archivo, fs, usuario)
    else:
        logger.warning(f"⚠️ Formato no soportado: {ext}")
        return

    if not unidades_contenido:
        logger.warning("⚠️ No se extrajo contenido.")
        return

    # Crear Documento Maestro con el campo de usuario
    documento = {
        "usuario_propietario": usuario,  # <--- LLAVE DE USUARIO
        "nombre_archivo": nombre,
        "tipo_archivo": ext.replace('.', ''),
        "fecha_ingesta": datetime.datetime.utcnow(),
        "total_unidades": len(unidades_contenido),
        "unidades_contenido": unidades_contenido,
        "estado_procesamiento": "PENDIENTE",
        "metadata": {
            "tamano_bytes": os.path.getsize(ruta_archivo),
            "version_modelo": "v3.0_paginado"
        }
    }

    try:
        # Usamos update_one con upsert para evitar duplicados del mismo archivo por el mismo usuario
        filtro = {"nombre_archivo": nombre, "usuario_propietario": usuario}
        res = collection.replace_one(filtro, documento, upsert=True)
        
        accion = "Actualizado" if res.matched_count > 0 else "Creado"
        logger.info(f"✅ {accion} exitosamente en Atlas para {usuario}. ({len(unidades_contenido)} unidades).")
        
    except Exception as e:
        logger.error(f"❌ Error guardando en MongoDB: {e}")

# --- 5. EJECUCIÓN ---

if __name__ == "__main__":
    # 1. Pedir Usuario PRIMERO
    usuario_actual = pedir_usuario_gui()

    if not usuario_actual or usuario_actual.strip() == "":
        logger.error("❌ El nombre de usuario es obligatorio. Saliendo.")
    else:
        logger.info(f"👤 Bienvenido, {usuario_actual}.")
        
        # 2. Conectar BD
        col, fs = conectar_bd()
        
        if col is not None:
            # 3. Seleccionar Archivos
            base_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
            raw_dir = os.path.join(base_dir, "data", "raw")
            
            logger.info("📂 Abriendo selector de archivos...")
            archivos_seleccionados = seleccionar_archivos_gui(initial_dir=raw_dir)

            if not archivos_seleccionados:
                logger.warning("⚠️ No se seleccionaron archivos.")
            else:
                # 4. Procesar cada archivo con el usuario
                for ruta in archivos_seleccionados:
                    ingestar_archivo(ruta, col, fs, usuario_actual)
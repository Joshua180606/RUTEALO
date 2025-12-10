import os
import datetime
import gridfs
import pypdf
import logging
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from werkzeug.utils import secure_filename
from src.config import DB_NAME, COLS, RAW_DIR, get_genai_model
from src.database import get_database
from src.utils import retry, validate_exam_responses, validate_exam_structure

# Importaciones para IA y lógica de negocio
import pandas as pd
import json
import re

logger = logging.getLogger(__name__)

# --- CONFIGURACIÓN GENERATIVA (Centralizada) ---
model = get_genai_model()

# Constantes de Lógica Educativa
JERARQUIA_BLOOM = ["Recordar", "Comprender", "Aplicar", "Analizar", "Evaluar", "Crear"]
COL_EXAM_INI = "examen_inicial"
COL_RUTAS = "rutas_aprendizaje"
COL_RAW = "materiales_crudos"


# --- CONEXIÓN BD ---
def get_db(db_name: str = DB_NAME):
    """Get database using centralized connection."""
    return get_database(db_name)


# --- LÓGICA DE INGESTA (Existente) ---
def guardar_imagen_gridfs(fs, img_bytes, nombre_archivo, pagina_idx, img_idx, ext, usuario):
    filename = f"{usuario}_{nombre_archivo}_P{pagina_idx}_IMG{img_idx}{ext}"
    file_id = fs.put(
        img_bytes,
        filename=filename,
        metadata={"documento_padre": nombre_archivo, "pagina_origen": pagina_idx, "usuario_propietario": usuario},
    )
    return {"gridfs_id": file_id, "nombre_archivo": filename}


def procesar_archivo_web(ruta_archivo, usuario, db):
    """Procesa un archivo subido y lo guarda en MongoDB."""
    fs = gridfs.GridFS(db)
    collection = db[COLS["RAW"]]

    nombre = os.path.basename(ruta_archivo)
    ext = os.path.splitext(nombre)[1].lower()
    unidades_contenido = []

    logger.info(f"🌐 Procesando web: {nombre} para {usuario}")

    try:
        if ext == ".pdf":
            reader = pypdf.PdfReader(ruta_archivo)
            for i, page in enumerate(reader.pages):
                texto = page.extract_text() or ""
                unidades_contenido.append(
                    {
                        "indice": i + 1,
                        "tipo_unidad": "pagina",
                        "contenido_texto": texto,
                        "imagenes": [],
                        "metadata_bloom": None,
                    }
                )

        elif ext == ".docx":
            doc = Document(ruta_archivo)
            texto = "\n".join([p.text for p in doc.paragraphs])
            unidades_contenido.append(
                {
                    "indice": 1,
                    "tipo_unidad": "documento_completo",
                    "contenido_texto": texto,
                    "imagenes": [],
                    "metadata_bloom": None,
                }
            )

        elif ext == ".pptx":
            prs = Presentation(ruta_archivo)
            for i, slide in enumerate(prs.slides):
                texto = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texto += shape.text + "\n"
                unidades_contenido.append(
                    {
                        "indice": i + 1,
                        "tipo_unidad": "diapositiva",
                        "contenido_texto": texto,
                        "imagenes": [],
                        "metadata_bloom": None,
                    }
                )

        if unidades_contenido:
            doc_data = {
                "usuario_propietario": usuario,
                "nombre_archivo": nombre,
                "tipo_archivo": ext.replace(".", ""),
                "fecha_ingesta": datetime.datetime.utcnow(),
                "unidades_contenido": unidades_contenido,
                "estado_procesamiento": "PENDIENTE",
            }
            collection.replace_one({"nombre_archivo": nombre, "usuario_propietario": usuario}, doc_data, upsert=True)
            return True, len(unidades_contenido)

    except Exception as e:
        logger.error(f"Error procesando archivo: {e}")
        return False, str(e)

    return False, "Formato no soportado o error desconocido"


# --- LÓGICA DE ETIQUETADO BLOOM (Automática) ---
def auto_etiquetar_bloom(usuario, db):
    """Busca documentos PENDIENTES del usuario y les aplica Bloom con Gemini."""
    col = db[COLS["RAW"]]
    docs = list(col.find({"usuario_propietario": usuario, "estado_procesamiento": "PENDIENTE"}))

    contexto_bloom = "Reglas de Bloom: Recordar, Comprender, Aplicar, Analizar, Evaluar, Crear."

    count = 0
    for doc in docs:
        unidades = doc.get("unidades_contenido", [])
        unidades_updated = []

        for u in unidades:
            texto = u.get("contenido_texto", "")[:1000]
            if not texto.strip():
                u["Categoria_Bloom"] = "Otro"
                u["Pedagogia_Detalle"] = {"justificacion": "Sin texto"}
                unidades_updated.append(u)
                continue

            prompt = f"""
            Clasifica este texto educativo según la Taxonomía de Bloom.
            Texto: {texto}
            Reglas: {contexto_bloom}
            Responde SOLO JSON: {{"Categoria_Bloom": "Nivel", "Justificacion": "Breve"}}
            Si no es educativo, categoria "Otro".
            """

            try:
                res = model.generate_content(prompt)
                json_res = json.loads(re.sub(r"```json|```", "", res.text).strip())
                u["Categoria_Bloom"] = json_res.get("Categoria_Bloom", "Otro")
                u["Pedagogia_Detalle"] = {"justificacion": json_res.get("Justificacion", "")}
            except Exception as e:
                logger.error(f"Error tagging with Bloom: {str(e)}")
                u["Categoria_Bloom"] = "Otro"
                u["Pedagogia_Detalle"] = {"error": "Fallo IA"}

            unidades_updated.append(u)

        col.update_one(
            {"_id": doc["_id"]},
            {"$set": {"unidades_contenido": unidades_updated, "estado_procesamiento": "BLOOM_COMPLETADO"}},
        )
        count += 1

    return count


# --- NUEVA LÓGICA: GENERADOR DE RUTAS (MOTOR PROMPTING ADAPTADO) ---


def obtener_contexto_usuario(db, usuario):
    """Recopila todo el texto procesado de este usuario, agrupado por categoría Bloom."""
    # Buscamos en la colección RAW usando la constante definida o importada
    col_raw = db[COLS["RAW"]]
    docs = col_raw.find({"usuario_propietario": usuario, "estado_procesamiento": "BLOOM_COMPLETADO"})

    contenido_por_nivel = {nivel: [] for nivel in JERARQUIA_BLOOM}
    contenido_total = ""

    for doc in docs:
        for unidad in doc.get("unidades_contenido", []):
            cat = unidad.get("Categoria_Bloom", "Otro")
            texto = unidad.get("contenido_texto", "")

            # Mapeo simple por si la IA usó sinónimos o mayúsculas
            for nivel in JERARQUIA_BLOOM:
                if nivel.lower() in cat.lower():
                    contenido_por_nivel[nivel].append(texto)
                    contenido_total += texto + "\n"
                    break

    return contenido_por_nivel, contenido_total


def cargar_marcos_pedagogicos():
    """Carga los CSV con teoría pedagógica (ZDP, Bloom, Flow) para guiar la generación del examen."""
    import pandas as pd
    import os
    
    base_path = os.path.join(os.path.dirname(__file__), "..", "data", "processed")
    
    marcos = {
        "bloom": None,
        "zdp": None,
        "flow": None
    }
    
    try:
        # Cargar df_bloom.csv
        bloom_path = os.path.join(base_path, "df_bloom.csv")
        if os.path.exists(bloom_path):
            marcos["bloom"] = pd.read_csv(bloom_path)
            logger.info(f"✅ Cargado df_bloom.csv: {len(marcos['bloom'])} filas")
        
        # Cargar df_zdp.csv
        zdp_path = os.path.join(base_path, "df_zdp.csv")
        if os.path.exists(zdp_path):
            marcos["zdp"] = pd.read_csv(zdp_path)
            logger.info(f"✅ Cargado df_zdp.csv: {len(marcos['zdp'])} filas")
        
        # Cargar df_flow.csv
        flow_path = os.path.join(base_path, "df_flow.csv")
        if os.path.exists(flow_path):
            marcos["flow"] = pd.read_csv(flow_path)
            logger.info(f"✅ Cargado df_flow.csv: {len(marcos['flow'])} filas")
    
    except Exception as e:
        logger.warning(f"⚠️ Error cargando marcos pedagógicos: {e}")
    
    return marcos


@retry(max_attempts=3, delay=2.0, backoff=2.0, exceptions=(Exception,))
def generar_examen_inicial(contenido_total):
    """Genera un examen diagnóstico CON PREGUNTAS REALES SOBRE EL MATERIAL del usuario.
    
    NUEVA ESTRATEGIA (Diciembre 2025):
    - Usa marcos pedagógicos de minería de datos (CSV: df_zdp, df_bloom, df_flow)
    - Genera preguntas SOBRE CONCEPTOS DEL MATERIAL (no autoevaluación)
    - Según las respuestas correctas/incorrectas, determina el nivel Bloom del estudiante
    - Incluye opción "e) No lo sé / Omitir" obligatoria
    
    Se reintenta automáticamente si falla.
    """
    if not contenido_total:
        return {}

    # Cargar marcos pedagógicos desde CSV
    marcos = cargar_marcos_pedagogicos()
    
    # Construir contexto pedagógico desde CSV para guiar la IA
    contexto_pedagogico = "\n📚 MARCOS TEÓRICOS PARA EVALUACIÓN PEDAGÓGICA:\n"
    
    # Marco Bloom (procesos cognitivos y tipos de conocimiento)
    if marcos["bloom"] is not None and len(marcos["bloom"]) > 0:
        contexto_pedagogico += "\n🔷 TAXONOMÍA DE BLOOM (Procesos cognitivos):\n"
        for _, row in marcos["bloom"].iterrows():
            cat = row.get("cat_bloom", "N/A")
            desc = str(row.get("proc_desc", ""))[:200]  # Primeros 200 chars
            if desc and desc != "N/A" and desc != "nan":
                contexto_pedagogico += f"  • {cat}: {desc}\n"
    
    # Marco ZDP (principios del aprendizaje desarrollador)
    if marcos["zdp"] is not None and len(marcos["zdp"]) > 0:
        contexto_pedagogico += "\n🎯 ZONA DE DESARROLLO PRÓXIMO (Principios):\n"
        zdp_sample = marcos["zdp"].head(6)  # Primeros 6 principios
        for _, row in zdp_sample.iterrows():
            principio = row.get("principio_zdp", "N/A")
            bloom_sug = row.get("cat_bloom_sugerida", "N/A")
            if principio != "N/A":
                contexto_pedagogico += f"  • {principio} → Evaluar con nivel: {bloom_sug}\n"
    
    # Marco Flow (dimensiones de experiencia óptima)
    if marcos["flow"] is not None and len(marcos["flow"]) > 0:
        contexto_pedagogico += "\n⚡ TEORÍA DEL FLOW (Dimensiones de motivación):\n"
        flow_sample = marcos["flow"].head(4)
        for _, row in flow_sample.iterrows():
            dimension = row.get("dimension", "N/A")
            bloom_sug = row.get("cat_bloom", "N/A")
            if dimension != "N/A":
                contexto_pedagogico += f"  • {dimension} → Nivel: {bloom_sug}\n"

    prompt = f"""
    Eres un profesor universitario experto en evaluación pedagógica con doctorado en Ciencias de la Educación.
    
    TAREA: Genera un EXAMEN DIAGNÓSTICO para evaluar qué nivel de la Taxonomía de Bloom domina el estudiante sobre el material.
    
    {contexto_pedagogico}
    
    📄 MATERIAL DEL ESTUDIANTE (contenido que subió):
    {contenido_total[:15000]}
    
    ⚠️ IMPORTANTE: Debes hacer preguntas SOBRE EL CONTENIDO ESPECÍFICO del material, NO preguntas meta-cognitivas.
    
    ❌ MAL (preguntas genéricas que NO evalúan conocimiento real):
       - "¿Cuál es el concepto central que identificas en el material?"
       - "¿Qué objetivo principal se persigue en el material?"
       - "¿Qué limitación identificas en tu comprensión actual?"
    
    ✅ BIEN (preguntas sobre conceptos, definiciones y temas del material):
       - "¿Qué es una base de datos relacional según el texto?" (si el tema es BD)
       - "¿Cuál es la diferencia entre normalización 2NF y 3NF mencionada?" (si el tema es BD)
       - "Según el material, ¿qué ventaja tiene el índice B-tree sobre el hash?" (si el tema es BD)
       - "¿Cómo implementarías una transacción ACID en PostgreSQL según lo explicado?" (si el tema es BD)
    
    REGLAS CRÍTICAS:
    1. Lee TODO el material y extrae CONCEPTOS CLAVE, DEFINICIONES, PROCESOS Y TEORÍAS mencionados
    
    2. Genera 10-12 preguntas distribuidas así:
       • 2-3 preguntas nivel RECORDAR: Definiciones, términos, hechos ("¿Qué es X?", "¿Cuál es la fórmula de Y?")
       • 2-3 preguntas nivel COMPRENDER: Explicaciones, interpretaciones ("¿Por qué ocurre X?", "Explica el concepto Y")
       • 2-3 preguntas nivel APLICAR: Casos prácticos ("¿Cómo usarías X para resolver Y?")
       • 2 preguntas nivel ANALIZAR: Comparaciones, relaciones ("¿Qué diferencia hay entre X e Y?")
       • 1-2 preguntas nivel EVALUAR: Juicios, críticas ("¿Cuál enfoque es mejor según...?")
       • 0-1 pregunta nivel CREAR: Diseño, síntesis ("¿Cómo combinarías X e Y?")
    
    3. TODAS las preguntas deben tener EXACTAMENTE 5 opciones: a, b, c, d + "e) No lo sé / Omitir"
    
    4. Las opciones incorrectas (distractores) deben:
       • Ser PLAUSIBLES (conceptos relacionados del mismo dominio)
       • Basarse en ERRORES COMUNES o conceptos similares del material
       • NO ser obviamente falsas
    
    5. La opción correcta debe estar TEXTUALMENTE o CONCEPTUALMENTE en el material (no inventar información)
    
    6. Si el material menciona ejemplos, úsalos en las preguntas de nivel Aplicar/Analizar
    
    FORMATO JSON OBLIGATORIO:
    {{
        "EXAMENES": {{
            "EXAMEN_INICIAL": [
                {{
                    "id": 1,
                    "pregunta": "Según el material, ¿qué es [CONCEPTO CLAVE del texto]?",
                    "opciones": [
                        "a) [Definición correcta extraída del material]",
                        "b) [Definición plausible pero incorrecta de un concepto relacionado]",
                        "c) [Otra definición incorrecta común]",
                        "d) [Distractor basado en error conceptual típico]",
                        "e) No lo sé / Omitir"
                    ],
                    "respuesta_correcta": "a",
                    "nivel_bloom_evaluado": "Recordar"
                }},
                {{
                    "id": 2,
                    "pregunta": "¿Por qué [PROCESO O CONCEPTO mencionado] funciona de esa manera según el texto?",
                    "opciones": [
                        "a) [Explicación plausible pero incorrecta]",
                        "b) [Explicación correcta del material]",
                        "c) [Otra explicación incorrecta]",
                        "d) [Confusión común sobre el tema]",
                        "e) No lo sé / Omitir"
                    ],
                    "respuesta_correcta": "b",
                    "nivel_bloom_evaluado": "Comprender"
                }},
                {{
                    "id": 3,
                    "pregunta": "Si quisieras implementar [TÉCNICA DEL MATERIAL], ¿cuál sería el primer paso según lo explicado?",
                    "opciones": [
                        "a) [Paso correcto del material]",
                        "b) [Paso plausible pero en orden incorrecto]",
                        "c) [Paso de otro proceso similar]",
                        "d) [Error común de implementación]",
                        "e) No lo sé / Omitir"
                    ],
                    "respuesta_correcta": "a",
                    "nivel_bloom_evaluado": "Aplicar"
                }}
            ]
        }}
    }}
    
    EJEMPLOS CONCRETOS SEGÚN EL DOMINIO DEL MATERIAL:
    
    Si el tema es BASES DE DATOS:
    - Recordar: "¿Qué propiedad garantiza que una transacción se ejecuta completamente o no se ejecuta?"
    - Comprender: "¿Por qué la normalización 3NF elimina dependencias transitivas?"
    - Aplicar: "¿Qué tipo de índice usarías para búsquedas de rango en una columna fecha?"
    - Analizar: "¿Cuál es la diferencia entre INNER JOIN y LEFT JOIN en SQL?"
    
    Si el tema es MACHINE LEARNING:
    - Recordar: "¿Qué es el overfitting según el material?"
    - Comprender: "¿Por qué el gradient descent puede quedar atrapado en mínimos locales?"
    - Aplicar: "¿Qué hiperparámetro ajustarías para reducir el overfitting en un árbol de decisión?"
    - Analizar: "Compara las ventajas de SVM vs Random Forest para clasificación binaria"
    
    Si el tema es ARQUITECTURA DE SOFTWARE:
    - Recordar: "¿Qué patrón de diseño resuelve el problema de crear familias de objetos relacionados?"
    - Comprender: "Explica por qué el patrón MVC separa la lógica de presentación"
    - Aplicar: "¿Cómo implementarías un Singleton thread-safe en Java?"
    - Analizar: "¿En qué se diferencia el patrón Strategy del patrón State?"
    
    CRÍTICO: 
    - Las preguntas deben permitir INFERIR el nivel Bloom del estudiante según:
      • Respuestas correctas en "Recordar" → Domina hechos y definiciones básicas
      • Respuestas correctas en "Comprender" → Entiende relaciones causa-efecto
      • Respuestas correctas en "Aplicar" → Puede usar conocimientos en casos prácticos
      • Respuestas correctas en "Analizar" → Puede descomponer y comparar conceptos
      • Respuestas correctas en "Evaluar" → Puede justificar decisiones
      • Respuestas correctas en "Crear" → Puede diseñar soluciones nuevas
    
    - El sistema usará el patrón de aciertos/fallos para determinar la Zona de Desarrollo Próximo (ZDP)
    - Si el estudiante responde "e) No lo sé / Omitir" se considera como "no domina ese nivel"
    """
    
    try:
        res = model.generate_content(prompt)
        text_clean = re.sub(r"```json|```", "", res.text).strip()
        data = json.loads(text_clean)
        
        # Validar que todas las preguntas tengan 5 opciones con la opción de omitir
        if "EXAMENES" in data and "EXAMEN_INICIAL" in data["EXAMENES"]:
            for pregunta in data["EXAMENES"]["EXAMEN_INICIAL"]:
                if len(pregunta.get("opciones", [])) < 5:
                    # Agregar opción omitir si falta
                    pregunta["opciones"].append("e) No lo sé / Omitir")
        
        return data
    except Exception as e:
        logger.error(f"⚠️ Error generando examen inicial: {e}")
        return {}


@retry(max_attempts=3, delay=2.0, backoff=2.0, exceptions=(Exception,))
def generar_bloque_ruta(nivel_bloom, textos_nivel, perfil_zdp=None, marcos=None):
    """Genera Flashcards y Exámenes para un nivel específico de Bloom usando funciones especializadas.
    
    Args:
        nivel_bloom (str): Nivel cognitivo a generar
        textos_nivel (list): Contenido del usuario para este nivel
        perfil_zdp (dict): Perfil ZDP del estudiante (opcional)
        marcos (dict): Marcos pedagógicos de CSV (opcional)
    
    Returns:
        dict: {"FLASHCARDS": [...], "EXAMENES": [...]} o None si debe omitirse
        
    FASE 2: Usa funciones especializadas con teoría pedagógica explícita.
    """
    if not textos_nivel:
        return None

    # Importar funciones especializadas
    from src.generadores_pedagogicos import generar_flashcards_con_teoria, generar_tests_con_teoria

    # Determinar estrategia según perfil ZDP
    estrategia = "estandar"
    
    if perfil_zdp:
        competentes = perfil_zdp.get("niveles_competentes", [])
        zona_proxima = perfil_zdp.get("zona_proxima", [])
        
        # OMITIR nivel si ya es competente
        if nivel_bloom in competentes:
            logger.info(f"⏭️  Omitiendo nivel {nivel_bloom} (ya competente)")
            return None
        
        # Aplicar scaffolding para zona próxima
        elif nivel_bloom in zona_proxima:
            estrategia = "scaffolding"
            logger.info(f"🎯 Generando {nivel_bloom} con SCAFFOLDING (zona próxima)")
        
        # Refuerzo intensivo para brechas
        else:
            estrategia = "refuerzo"
            logger.info(f"💪 Generando {nivel_bloom} con REFUERZO (brecha detectada)")
    else:
        logger.info(f"📝 Generando {nivel_bloom} con estrategia estándar (sin perfil ZDP)")

    # FASE 2: Generar flashcards y tests con funciones especializadas
    try:
        flashcards = generar_flashcards_con_teoria(
            nivel_bloom=nivel_bloom,
            textos_nivel=textos_nivel,
            estrategia=estrategia,
            marcos=marcos
        )
        
        tests = generar_tests_con_teoria(
            nivel_bloom=nivel_bloom,
            textos_nivel=textos_nivel,
            estrategia=estrategia,
            marcos=marcos
        )
        
        # Validar que se generó contenido
        if not flashcards and not tests:
            logger.warning(f"⚠️  No se generó contenido para {nivel_bloom}")
            return None
        
        return {
            "FLASHCARDS": flashcards,
            "EXAMENES": tests
        }
        
    except Exception as e:
        logger.error(f"❌ Error generando bloque {nivel_bloom}: {e}")
        return None


def generar_ruta_aprendizaje(usuario, db):
    """
    Orquestador principal: Lee todo el material del usuario y (re)genera la ruta completa.
    Retorna un mensaje de estado.
    """
    logger.info(f"🛤️ Iniciando generación de ruta para: {usuario}")

    # 1. Obtener Contexto Global
    contenido_bloom, contenido_total_raw = obtener_contexto_usuario(db, usuario)
    col_examen = db[COL_EXAM_INI]
    col_ruta = db[COL_RUTAS]

    if not contenido_total_raw:
        # Si no hay nuevo contenido Bloom, intentar reutilizar lo existente o crear un mínimo base
        examen_existente = col_examen.find_one({"usuario": usuario})
        ruta_existente = col_ruta.find_one({"usuario": usuario})

        if examen_existente or ruta_existente:
            logger.info("No hay nuevo contenido Bloom; se mantiene la ruta/examen previos.")
            return "No hay nuevo contenido Bloom; se mantuvo la ruta y examen existentes."

        logger.warning("No hay contenido Bloom; generando materiales base mínimos para no bloquear la ruta.")

        examen_minimo = _crear_examen_minimo()
        doc_examen_fallback = {
            "usuario": usuario,
            "contenido": examen_minimo,
            "estado": "PENDIENTE",
            "origen": "fallback_sin_bloom",
            "fecha_generacion": datetime.datetime.utcnow(),
        }
        col_examen.replace_one({"usuario": usuario}, doc_examen_fallback, upsert=True)

        ruta_fallback = _crear_ruta_minima(usuario)
        col_ruta.replace_one({"usuario": usuario}, ruta_fallback, upsert=True)

        return "Ruta generada con materiales base mínimos (sin Bloom). Carga más contenido para personalizarla."

    # 2. Generar/Actualizar Examen Inicial (ZDP)
    # Siempre regeneramos para incluir el nuevo material en el diagnóstico
    examen_ini_data = generar_examen_inicial(contenido_total_raw)

    if examen_ini_data:
        doc_examen_ini = {
            "usuario": usuario,
            "contenido": examen_ini_data,
            "estado": "PENDIENTE",
            "fecha_generacion": datetime.datetime.utcnow(),
        }
        # Guardamos en la colección correspondiente
        db[COL_EXAM_INI].replace_one({"usuario": usuario}, doc_examen_ini, upsert=True)

    # NUEVO: Obtener perfil ZDP (si existe evaluación previa)
    from src.models.evaluacion_zdp import EvaluadorZDP
    evaluador = EvaluadorZDP()
    perfil_zdp = evaluador.obtener_perfil_zdp_simple(usuario)
    
    if perfil_zdp:
        logger.info(f"📊 Usando perfil ZDP: Omitiendo {len(perfil_zdp['niveles_competentes'])} niveles competentes")
    else:
        logger.info("📝 Sin perfil ZDP previo, generando ruta completa")

    # NUEVO: Cargar marcos pedagógicos (CSV)
    marcos = cargar_marcos_pedagogicos()

    # 3. Generar Ruta de Aprendizaje ADAPTATIVA (Flow + ZDP)
    logger.info("🛤️ Diseñando Ruta de Aprendizaje Personalizada...")
    
    ruta_completa = {}
    secuencia_id = 1
    niveles_omitidos = []
    niveles_generados = []

    for nivel in JERARQUIA_BLOOM:
        textos = contenido_bloom.get(nivel, [])
        if not textos:
            continue

        # NUEVO: Pasar perfil_zdp y marcos a generar_bloque_ruta
        logger.debug(f"   ⚡ Procesando Nivel: {nivel}...")
        bloque_generado = generar_bloque_ruta(nivel, textos, perfil_zdp, marcos)

        if bloque_generado is None:
            # Nivel omitido (competente)
            ruta_completa[nivel] = {
                "id_orden": secuencia_id,
                "bloqueado": False,  # Desbloqueado pero omitido
                "estado": "OMITIDO",
                "razon": "El estudiante ya domina este nivel según evaluación ZDP",
                "contenido": {"FLASHCARDS": [], "EXAMENES": []}
            }
            niveles_omitidos.append(nivel)
            secuencia_id += 1
        elif bloque_generado:
            ruta_completa[nivel] = {
                "id_orden": secuencia_id,
                "bloqueado": True if secuencia_id > 1 else False,  # El primero desbloqueado
                "contenido": bloque_generado,
            }
            niveles_generados.append(nivel)
            secuencia_id += 1

    # 4. Guardar Estructura Final en Mongo
    doc_ruta = {
        "usuario": usuario,
        "estructura_ruta": {
            "usuario": usuario,
            "examenes": {nivel: data["contenido"].get("EXAMENES", []) for nivel, data in ruta_completa.items()},
            "flashcards": {nivel: data["contenido"].get("FLASHCARDS", []) for nivel, data in ruta_completa.items()},
        },
        "metadatos_ruta": {
            "niveles_incluidos": list(ruta_completa.keys()),
            "niveles_generados": niveles_generados,
            "niveles_omitidos": niveles_omitidos,
            "progreso_global": 0,
            "personalizada_zdp": perfil_zdp is not None,
            "estado_niveles": {
                nivel: data.get("estado", "BLOQUEADO" if data["bloqueado"] else "DISPONIBLE") 
                for nivel, data in ruta_completa.items()
            },
        },
        "fecha_actualizacion": datetime.datetime.utcnow(),
    }
    
    # Agregar info de perfil ZDP si existe
    if perfil_zdp:
        doc_ruta["metadatos_ruta"]["nivel_actual_estudiante"] = perfil_zdp.get("nivel_actual")
        doc_ruta["metadatos_ruta"]["zona_proxima"] = perfil_zdp.get("zona_proxima", [])

    # Si no se generaron bloques (p. ej. contenidos vacíos), usar un fallback mínimo
    if not ruta_completa:
        logger.warning("No se generaron bloques de ruta; creando ruta mínima para evitar bloqueo.")
        ruta_fallback = _crear_ruta_minima(usuario)
        col_ruta.replace_one({"usuario": usuario}, ruta_fallback, upsert=True)
        return "Ruta generada con materiales base mínimos. Agrega más contenido para enriquecerla."

    col_ruta.replace_one({"usuario": usuario}, doc_ruta, upsert=True)

    # NUEVO: Log con estadísticas de optimización
    if niveles_omitidos:
        logger.info(f"✅ Ruta OPTIMIZADA generada:")
        logger.info(f"   📊 Total niveles: {len(ruta_completa)}")
        logger.info(f"   ✅ Generados: {len(niveles_generados)} {niveles_generados}")
        logger.info(f"   ⏭️  Omitidos: {len(niveles_omitidos)} {niveles_omitidos}")
        logger.info(f"   💰 Ahorro estimado: ~{len(niveles_omitidos) * 40}% tokens")
        return f"Ruta PERSONALIZADA regenerada ({len(niveles_generados)} niveles activos, {len(niveles_omitidos)} omitidos por dominio)."
    else:
        logger.info(f"✅ Ruta completa generada con {len(ruta_completa)} niveles Bloom.")
        return f"Ruta regenerada con {len(ruta_completa)} niveles y Examen Diagnóstico actualizado."


def _crear_examen_minimo():
    """Crea un examen diagnóstico base cuando no hay contenido Bloom disponible.
    
    ACTUALIZADO: Ya no usa preguntas meta-cognitivas genéricas.
    Ahora pide al usuario que describa conceptos básicos del material que subió.
    """
    preguntas_base = [
        {
            "id": 1,
            "pregunta": "¿Cuál es el concepto técnico MÁS IMPORTANTE mencionado en tu material?",
            "opciones": [
                "a) [Describe brevemente el concepto principal que encontraste]",
                "b) No identifiqué ningún concepto técnico específico",
                "c) El material solo tiene información general",
                "d) No revisé el material completo",
                "e) No lo sé / Omitir",
            ],
            "respuesta_correcta": "a",
            "nivel_bloom_evaluado": "Recordar",
        },
        {
            "id": 2,
            "pregunta": "Según el material, ¿para QUÉ SIRVE o CÓMO FUNCIONA ese concepto principal?",
            "opciones": [
                "a) [Explica brevemente su función o propósito según el texto]",
                "b) No lo explica claramente",
                "c) Solo lo menciona sin explicarlo",
                "d) No entendí su función",
                "e) No lo sé / Omitir",
            ],
            "respuesta_correcta": "a",
            "nivel_bloom_evaluado": "Comprender",
        },
        {
            "id": 3,
            "pregunta": "Si tuvieras que USAR ese concepto en un proyecto real, ¿qué EJEMPLO concreto mencionó el material?",
            "opciones": [
                "a) [Describe el ejemplo o caso de uso que mencionó el texto]",
                "b) No mencionó ejemplos prácticos",
                "c) Los ejemplos no eran claros",
                "d) No recuerdo ejemplos específicos",
                "e) No lo sé / Omitir",
            ],
            "respuesta_correcta": "a",
            "nivel_bloom_evaluado": "Aplicar",
        },
        {
            "id": 4,
            "pregunta": "¿El material menciona DIFERENCIAS o COMPARACIONES entre ese concepto y otros relacionados?",
            "opciones": [
                "a) Sí, compara con [menciona el otro concepto comparado]",
                "b) No hace comparaciones",
                "c) Solo menciona conceptos aislados",
                "d) No presté atención a las comparaciones",
                "e) No lo sé / Omitir",
            ],
            "respuesta_correcta": "a",
            "nivel_bloom_evaluado": "Analizar",
        },
        {
            "id": 5,
            "pregunta": "Según TU criterio y lo leído, ¿cuál es la LIMITACIÓN o DESVENTAJA del concepto principal?",
            "opciones": [
                "a) [Describe una limitación mencionada o que identificaste]",
                "b) No tiene limitaciones según el texto",
                "c) El texto no analiza limitaciones",
                "d) No puedo identificar limitaciones aún",
                "e) No lo sé / Omitir",
            ],
            "respuesta_correcta": "a",
            "nivel_bloom_evaluado": "Evaluar",
        },
    ]

    return {"EXAMENES": {"EXAMEN_INICIAL": preguntas_base}}


def _crear_ruta_minima(usuario):
    """Genera una ruta mínima con flashcards y exámenes genéricos para no bloquear el flujo."""
    niveles_base = JERARQUIA_BLOOM[:3]
    flashcards = {}
    examenes = {}
    estado_niveles = {}

    for idx, nivel in enumerate(niveles_base, start=1):
        flashcards[nivel] = [
            {
                "id": 1,
                "frente": f"Idea clave de {nivel}",
                "reverso": "Resume con tus palabras la idea principal de tu material.",
                "visto": False,
            }
        ]
        examenes[nivel] = [
            {
                "id": 1,
                "pregunta": f"Aplica un concepto de {nivel} a tu experiencia diaria.",
                "opciones": [
                    "a) Relacionar con un ejemplo propio",
                    "b) Repetir sin aplicar",
                ],
                "respuesta_correcta": "a",
                "realizado": False,
            }
        ]
        estado_niveles[nivel] = "DISPONIBLE" if idx == 1 else "BLOQUEADO"

    return {
        "usuario": usuario,
        "estructura_ruta": {
            "usuario": usuario,
            "examenes": examenes,
            "flashcards": flashcards,
        },
        "metadatos_ruta": {
            "niveles_incluidos": niveles_base,
            "progreso_global": 0,
            "estado_niveles": estado_niveles,
        },
        "fecha_actualizacion": datetime.datetime.utcnow(),
        "origen": "fallback_sin_bloom",
    }


# --- NOTA: Funciones ZDP movidas a src/models/evaluacion_zdp.py ---
# Importar directamente desde allí:
# - evaluar_examen_simple() -> para procesar respuestas
# - obtener_perfil_zdp() -> para obtener perfil del estudiante
# - EvaluadorZDP.generar_ruta_personalizada() -> para rutas personalizadas


# --- FUNCIONES NUEVAS PARA REDISEÑO DASHBOARD ---


def procesar_multiples_archivos_web(archivos_rutas: list, usuario: str, db) -> tuple:
    """
    Procesa múltiples archivos en una operación.
    
    Args:
        archivos_rutas: List[str] - Rutas locales de archivos
        usuario: str - Usuario propietario
        db: Database - Instancia MongoDB
    
    Returns:
        Tuple[bool, List[dict], str] - (éxito, resultados por archivo, mensaje)
        
    Formato de resultados:
        [
            {"nombre": "documento.pdf", "unidades": 10, "estado": "OK", "error": None},
            {"nombre": "invalido.txt", "unidades": 0, "estado": "ERROR", "error": "..."}
        ]
    """
    resultados = []
    total_unidades = 0
    
    for ruta_archivo in archivos_rutas:
        try:
            # Validar que archivo existe
            if not os.path.exists(ruta_archivo):
                resultados.append({
                    "nombre": os.path.basename(ruta_archivo),
                    "unidades": 0,
                    "estado": "ERROR",
                    "error": "Archivo no encontrado"
                })
                continue
            
            # Procesar archivo
            ok, msg = procesar_archivo_web(ruta_archivo, usuario, db)
            
            if ok:
                # msg contiene cantidad de unidades procesadas
                resultados.append({
                    "nombre": os.path.basename(ruta_archivo),
                    "unidades": msg,
                    "estado": "OK",
                    "error": None
                })
                total_unidades += msg
            else:
                resultados.append({
                    "nombre": os.path.basename(ruta_archivo),
                    "unidades": 0,
                    "estado": "ERROR",
                    "error": msg
                })
        
        except Exception as e:
            logger.error(f"Error procesando {ruta_archivo}: {e}")
            resultados.append({
                "nombre": os.path.basename(ruta_archivo),
                "unidades": 0,
                "estado": "ERROR",
                "error": str(e)
            })
    
    # Determinar éxito global
    exitosos = [r for r in resultados if r["estado"] == "OK"]
    éxito = len(exitosos) > 0
    
    # Mensaje resumen
    msg_resumen = f"Procesados {len(exitosos)}/{len(archivos_rutas)} archivos ({total_unidades} unidades)"
    
    return éxito, resultados, msg_resumen


def obtener_rutas_usuario(usuario: str, db) -> list:
    """
    Obtiene lista de rutas del usuario con metadata.
    
    Args:
        usuario: str - ID del usuario
        db: Database - Instancia MongoDB
    
    Returns:
        List[dict] - Rutas ordenadas por fecha_actualizacion DESC
    """
    col = db[COLS["RUTAS"]]
    
    try:
        rutas_cursor = col.find(
            {"usuario": usuario},
            {
                "_id": 1,
                "nombre_ruta": 1,
                "descripcion": 1,
                "estado": 1,
                "progreso_global": 1,
                "fecha_actualizacion": 1,
                "archivos_fuente": 1,
                "metadatos_ruta.niveles_incluidos": 1
            }
        ).sort("fecha_actualizacion", -1)
        
        rutas_list = []
        for ruta in rutas_cursor:
            niveles = ruta.get("metadatos_ruta", {}).get("niveles_incluidos", [])
            archivos = ruta.get("archivos_fuente", [])
            
            rutas_list.append({
                "ruta_id": str(ruta["_id"]),
                "nombre_ruta": ruta.get("nombre_ruta", "Sin nombre"),
                "descripcion": ruta.get("descripcion", ""),
                "estado": ruta.get("estado", "ACTIVA"),
                "progreso": ruta.get("progreso_global", 0),
                "archivos_count": len(archivos),
                "niveles_completados": len(niveles),
                "fecha_actualizacion": ruta.get("fecha_actualizacion")
            })
        
        return rutas_list
    
    except Exception as e:
        logger.error(f"Error obteniendo rutas para {usuario}: {e}")
        return []
